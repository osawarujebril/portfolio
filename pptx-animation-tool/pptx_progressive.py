"""
Progressive Slide Builder v2.0
===============================
Builds progressive reveal sequences as separate slides — each slide adds one more element.
This is how Fladlien actually does it. Works in ANY presentation software.

v2.0 additions:
- Image placement (add_image, add_background_image)
- Rich text (multiple colors/fonts/sizes in one text box)
- Gradient backgrounds
- Dark overlays on background images
- Shape support (circles, lines, connectors, arrows)
- Multiple font families per slide

Usage:
    from tools.pptx_progressive import ProgressiveSlideBuilder, simple_slide

    builder = ProgressiveSlideBuilder(prs)

    # Progressive reveal
    builder.start_sequence(bg_color=(0, 0, 0))
    builder.add_persistent("3 REASONS WHY", x=1, y=0.5, w=14, h=1.5, size=52, bold=True, color=(255, 192, 0))
    builder.reveal("1.  The old playbook is dead", x=2, y=2.5, w=12, h=1.2, size=40)
    builder.reveal("2.  The winners changed the rules", x=2, y=4.2, w=12, h=1.2, size=40)
    builder.reveal("3.  You haven't seen this yet...", x=2, y=5.9, w=12, h=1.2, size=40, color=(5, 203, 2))
    builder.build()
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree


# ─────────────────────────────────────────────
# DESIGN TOKENS (Fladlien Color-as-Meaning System)
# ─────────────────────────────────────────────
COLORS = {
    'white':      (255, 255, 255),
    'gold':       (255, 192, 0),
    'green':      (5, 203, 2),
    'red':        (226, 33, 70),
    'yellow':     (255, 255, 0),
    'bright_red': (255, 0, 0),
    'blue':       (0, 0, 255),
    'black':      (0, 0, 0),
    'dark':       (10, 10, 20),
    'golden':     (249, 185, 0),
}

ALIGN_MAP = {
    "center": PP_ALIGN.CENTER,
    "left":   PP_ALIGN.LEFT,
    "right":  PP_ALIGN.RIGHT,
}

A_NS = '{http://schemas.openxmlformats.org/drawingml/2006/main}'


def _set_shape_alpha(shape, opacity):
    """
    Set transparency on a shape's solid fill via XML injection.

    Args:
        shape: python-pptx shape with solid fill already applied
        opacity: 0.0 (fully transparent) to 1.0 (fully opaque)
    """
    alpha_val = str(int((1 - opacity) * 100000))
    # Navigate the shape's XML to find the srgbClr element
    sp_elem = shape._element
    srgb_elems = sp_elem.findall('.//' + A_NS + 'srgbClr')
    # Use the first srgbClr in solidFill (not line or other fills)
    for srgb in srgb_elems:
        parent = srgb.getparent()
        if parent is not None and parent.tag.endswith('solidFill'):
            alpha_elem = etree.SubElement(srgb, A_NS + 'alpha')
            alpha_elem.set('val', alpha_val)
            break


def _create_overlay(slide, r, g, b, opacity, width, height):
    """
    Create a semi-transparent color overlay rectangle.

    Args:
        slide: python-pptx slide
        r, g, b: color values 0-255
        opacity: 0.0-1.0
        width, height: slide dimensions (Emu or Inches)
    Returns:
        The overlay shape object
    """
    overlay_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        width, height
    )
    overlay_shape.fill.solid()
    overlay_shape.fill.fore_color.rgb = RGBColor(r, g, b)
    overlay_shape.line.fill.background()  # No border
    _set_shape_alpha(overlay_shape, opacity)
    return overlay_shape


class ProgressiveSlideBuilder:
    """Builds progressive reveal sequences as separate slides."""

    def __init__(self, prs, slide_width=16, slide_height=9):
        self.prs = prs
        self.slide_width = slide_width
        self.slide_height = slide_height
        self._persistent_elements = []
        self._reveal_elements = []
        self._bg_color = (0, 0, 0)
        self._bg_gradient = None
        self._bg_image = None
        self._bg_overlay = None
        self._current_notes = None

    def start_sequence(self, bg_color=(0, 0, 0), bg_gradient=None,
                       bg_image=None, bg_overlay=None):
        """
        Start a new progressive sequence. Resets all elements.

        Args:
            bg_color: tuple (r, g, b) — solid background color
            bg_gradient: dict with 'angle', 'stops' — gradient background
                Example: {'angle': 270, 'stops': [(0, (0,0,0)), (1, (26,26,26))]}
            bg_image: str path to image file — background image
            bg_overlay: tuple (r, g, b, opacity) — overlay on background image
                Example: (0, 0, 0, 0.6) for 60% black overlay
        """
        self._persistent_elements = []
        self._reveal_elements = []
        self._bg_color = bg_color
        self._bg_gradient = bg_gradient
        self._bg_image = bg_image
        self._bg_overlay = bg_overlay
        self._current_notes = None

    # ─────────────────────────────────────────
    # PERSISTENT ELEMENTS (appear on every slide)
    # ─────────────────────────────────────────

    def add_persistent(self, text, x, y, w, h, size=44, bold=False,
                       color=(255, 255, 255), align="center", font_name=None):
        """Add a text element that appears on EVERY slide in the sequence."""
        self._persistent_elements.append({
            "type": "text",
            "text": text, "x": x, "y": y, "w": w, "h": h,
            "size": size, "bold": bold, "color": color, "align": align,
            "font_name": font_name,
        })

    def add_persistent_rich_text(self, runs, x, y, w, h, align="center"):
        """
        Add a rich text element (multiple fonts/colors in one box) on every slide.

        Args:
            runs: list of dicts, each with:
                - text: str
                - size: int (pt)
                - color: tuple (r,g,b)
                - bold: bool (optional, default False)
                - font_name: str (optional)
                - italic: bool (optional)
            x, y, w, h: position and size in inches
            align: "left" / "center" / "right"
        """
        self._persistent_elements.append({
            "type": "rich_text",
            "runs": runs, "x": x, "y": y, "w": w, "h": h,
            "align": align,
        })

    def add_persistent_shape(self, shape_type, x, y, w, h,
                             fill_color=(255, 192, 0), line_color=None,
                             line_width=None):
        """
        Add a shape that appears on every slide.

        Args:
            shape_type: MSO_SHAPE constant or string name
                Strings: 'rectangle', 'oval', 'circle', 'line',
                         'rounded_rectangle', 'arrow_right', 'chevron',
                         'triangle', 'diamond', 'pentagon', 'hexagon',
                         'funnel'
            fill_color: tuple (r,g,b) or None for no fill
            line_color: tuple (r,g,b) or None for no line
            line_width: float in pt or None
        """
        self._persistent_elements.append({
            "type": "shape",
            "shape_type": shape_type, "x": x, "y": y, "w": w, "h": h,
            "fill_color": fill_color, "line_color": line_color,
            "line_width": line_width,
        })

    def add_persistent_image(self, image_path, x, y, w, h):
        """Add an image that appears on EVERY slide in the sequence."""
        self._persistent_elements.append({
            "type": "image",
            "image_path": image_path, "x": x, "y": y, "w": w, "h": h,
        })

    # ─────────────────────────────────────────
    # REVEAL ELEMENTS (each creates a new slide)
    # ─────────────────────────────────────────

    def reveal(self, text, x, y, w, h, size=40, bold=False,
               color=(255, 255, 255), align="left", notes=None,
               font_name=None):
        """Add a text element that gets REVEALED (creates a new slide)."""
        self._reveal_elements.append({
            "type": "text",
            "text": text, "x": x, "y": y, "w": w, "h": h,
            "size": size, "bold": bold, "color": color, "align": align,
            "font_name": font_name,
            "notes": notes,
        })

    def reveal_rich_text(self, runs, x, y, w, h, align="left", notes=None):
        """
        Reveal a rich text element (multiple fonts/colors).

        Args:
            runs: list of dicts — same format as add_persistent_rich_text
        """
        self._reveal_elements.append({
            "type": "rich_text",
            "runs": runs, "x": x, "y": y, "w": w, "h": h,
            "align": align,
            "notes": notes,
        })

    def reveal_image(self, image_path, x, y, w, h, notes=None):
        """Reveal an image (creates a new slide with the image added)."""
        self._reveal_elements.append({
            "type": "image",
            "image_path": image_path, "x": x, "y": y, "w": w, "h": h,
            "notes": notes,
        })

    def reveal_shape(self, shape_type, x, y, w, h, fill_color=(255, 192, 0),
                     line_color=None, line_width=None, notes=None):
        """Reveal a shape element."""
        self._reveal_elements.append({
            "type": "shape",
            "shape_type": shape_type, "x": x, "y": y, "w": w, "h": h,
            "fill_color": fill_color, "line_color": line_color,
            "line_width": line_width,
            "notes": notes,
        })

    def reveal_group(self, elements, notes=None):
        """
        Reveal multiple elements on the SAME click (same slide transition).

        Args:
            elements: list of dicts, each with:
                For text: type="text", text, x, y, w, h, size, bold, color, align, font_name
                For rich_text: type="rich_text", runs, x, y, w, h, align
                For image: type="image", image_path, x, y, w, h
                For shape: type="shape", shape_type, x, y, w, h, fill_color
        """
        self._reveal_elements.append({
            "type": "group",
            "elements": elements,
            "notes": notes,
        })

    # ─────────────────────────────────────────
    # BUILD
    # ─────────────────────────────────────────

    def build(self):
        """Build all slides for this sequence. Returns list of slide objects."""
        slides = []

        if self._reveal_elements:
            # First slide: persistent elements only
            slide = self._make_slide(self._persistent_elements, [])
            slides.append(slide)

            # Each subsequent slide adds one more reveal element
            for i in range(len(self._reveal_elements)):
                revealed_so_far = self._reveal_elements[:i + 1]
                notes = self._reveal_elements[i].get("notes")
                slide = self._make_slide(self._persistent_elements, revealed_so_far, notes)
                slides.append(slide)
        else:
            slide = self._make_slide(self._persistent_elements, [])
            slides.append(slide)

        return slides

    # ─────────────────────────────────────────
    # INTERNAL: Slide construction
    # ─────────────────────────────────────────

    def _make_slide(self, persistent, revealed, notes=None):
        """Create a single slide with the given elements."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank

        # Background
        self._apply_background(slide)

        # Background image (if set)
        if self._bg_image:
            self._add_background_image(slide, self._bg_image)
            if self._bg_overlay:
                self._add_overlay(slide, self._bg_overlay)

        # Add persistent elements
        for elem in persistent:
            self._add_element(slide, elem)

        # Add revealed elements
        for elem in revealed:
            if elem["type"] == "group":
                for sub in elem["elements"]:
                    self._add_element(slide, sub)
            else:
                self._add_element(slide, elem)

        # Speaker notes
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

        return slide

    def _apply_background(self, slide):
        """Apply background color or gradient to a slide."""
        bg = slide.background
        fill = bg.fill

        if self._bg_gradient:
            fill.gradient()
            angle = self._bg_gradient.get('angle', 270)
            stops = self._bg_gradient.get('stops', [])

            # Set gradient stops
            for i, (position, color) in enumerate(stops):
                if i < len(fill.gradient_stops):
                    stop = fill.gradient_stops[i]
                    stop.color.rgb = RGBColor(*color)
                    stop.position = position

            # Set gradient angle via XML
            # The gradient XML is under the fill element's XML tree
            a_ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
            grad_fill_elem = fill._fill._element if hasattr(fill._fill, '_element') else fill._fill
            # Find or create the lin element for angle
            try:
                xml_elem = bg._element
                grad_fill_xml = xml_elem.find('.//' + a_ns + 'gradFill')
                if grad_fill_xml is not None:
                    lin_elem = grad_fill_xml.find(a_ns + 'lin')
                    if lin_elem is None:
                        lin_elem = etree.SubElement(grad_fill_xml, a_ns + 'lin')
                    lin_elem.set('ang', str(angle * 60000))
                    lin_elem.set('scaled', '1')
            except Exception:
                pass  # Gradient works without angle, just uses default direction
        else:
            fill.solid()
            fill.fore_color.rgb = RGBColor(*self._bg_color)

    def _add_background_image(self, slide, image_path):
        """
        Add a full-slide background image and move it to the back layer.
        The image fills the entire slide area.
        """
        pic = slide.shapes.add_picture(
            image_path,
            Inches(0), Inches(0),
            self.prs.slide_width, self.prs.slide_height
        )
        # Move picture to the back of the shape tree
        sp_tree = slide.shapes._spTree
        sp_tree.remove(pic._element)
        sp_tree.insert(2, pic._element)
        return pic

    def _add_overlay(self, slide, overlay):
        """
        Add a semi-transparent color overlay on top of a background image.

        Args:
            overlay: tuple (r, g, b, opacity) where opacity is 0.0-1.0
        """
        r, g, b, opacity = overlay
        overlay_shape = _create_overlay(slide, r, g, b, opacity,
                                        self.prs.slide_width, self.prs.slide_height)
        return overlay_shape

    def _add_element(self, slide, elem):
        """Add a single element to a slide."""
        elem_type = elem.get("type", "text")

        if elem_type == "text":
            return self._add_text(slide, elem)
        elif elem_type == "rich_text":
            return self._add_rich_text(slide, elem)
        elif elem_type == "image":
            return self._add_image(slide, elem)
        elif elem_type == "shape":
            return self._add_shape(slide, elem)

    def _add_text(self, slide, elem):
        """Add a simple text element."""
        box = slide.shapes.add_textbox(
            Inches(elem["x"]), Inches(elem["y"]),
            Inches(elem["w"]), Inches(elem["h"])
        )
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = elem["text"]
        run.font.size = Pt(elem["size"])
        run.font.bold = elem.get("bold", False)
        run.font.color.rgb = RGBColor(*elem["color"])

        if elem.get("font_name"):
            run.font.name = elem["font_name"]

        p.alignment = ALIGN_MAP.get(elem.get("align", "center"), PP_ALIGN.CENTER)
        return box

    def _add_rich_text(self, slide, elem):
        """
        Add a rich text element with multiple runs (different fonts/colors/sizes).

        Example runs:
        [
            {"text": "WHOLESALE", "size": 52, "color": (255,192,0), "bold": True, "font_name": "Impact"},
            {"text": " is the future", "size": 52, "color": (255,255,255), "font_name": "Arial"},
        ]
        """
        box = slide.shapes.add_textbox(
            Inches(elem["x"]), Inches(elem["y"]),
            Inches(elem["w"]), Inches(elem["h"])
        )
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]

        for run_spec in elem["runs"]:
            run = p.add_run()
            run.text = run_spec["text"]
            run.font.size = Pt(run_spec.get("size", 40))
            run.font.color.rgb = RGBColor(*run_spec.get("color", (255, 255, 255)))
            run.font.bold = run_spec.get("bold", False)
            run.font.italic = run_spec.get("italic", False)

            if run_spec.get("font_name"):
                run.font.name = run_spec["font_name"]

        p.alignment = ALIGN_MAP.get(elem.get("align", "center"), PP_ALIGN.CENTER)
        return box

    def _add_image(self, slide, elem):
        """Add an image to the slide."""
        pic = slide.shapes.add_picture(
            elem["image_path"],
            Inches(elem["x"]), Inches(elem["y"]),
            Inches(elem["w"]), Inches(elem["h"])
        )
        return pic

    def _add_shape(self, slide, elem):
        """Add a shape (rectangle, oval, line, arrow, etc.) to the slide."""
        shape_type = elem["shape_type"]

        # Convert string names to MSO_SHAPE constants
        if isinstance(shape_type, str):
            shape_map = {
                'rectangle':         MSO_SHAPE.RECTANGLE,
                'oval':              MSO_SHAPE.OVAL,
                'circle':            MSO_SHAPE.OVAL,
                'line':              MSO_SHAPE.RECTANGLE,  # thin rectangle as line
                'rounded_rectangle': MSO_SHAPE.ROUNDED_RECTANGLE,
                'arrow_right':       MSO_SHAPE.RIGHT_ARROW,
                'chevron':           MSO_SHAPE.CHEVRON,
                'triangle':          MSO_SHAPE.ISOSCELES_TRIANGLE,
                'diamond':           MSO_SHAPE.DIAMOND,
                'pentagon':          MSO_SHAPE.PENTAGON,
                'hexagon':           MSO_SHAPE.HEXAGON,
                'funnel':            MSO_SHAPE.FUNNEL,
            }
            shape_type = shape_map.get(shape_type.lower(), MSO_SHAPE.RECTANGLE)

        shape = slide.shapes.add_shape(
            shape_type,
            Inches(elem["x"]), Inches(elem["y"]),
            Inches(elem["w"]), Inches(elem["h"])
        )

        # Fill
        fill_color = elem.get("fill_color")
        if fill_color is not None:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*fill_color)
        else:
            shape.fill.background()  # No fill

        # Line/border
        line_color = elem.get("line_color")
        if line_color is not None:
            shape.line.color.rgb = RGBColor(*line_color)
            if elem.get("line_width"):
                shape.line.width = Pt(elem["line_width"])
        else:
            shape.line.fill.background()  # No line

        return shape


# ─────────────────────────────────────────────
# STANDALONE SLIDE FUNCTIONS
# ─────────────────────────────────────────────

def simple_slide(prs, text, size=96, bold=True, color=(255, 255, 255),
                 bg_color=(0, 0, 0), x=1, y=3, w=14, h=3, align="center",
                 notes=None, font_name=None):
    """
    Create a single simple slide (no animation needed).
    Good for: single-word impacts, questions, transitions.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*bg_color)

    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)

    if font_name:
        run.font.name = font_name

    p.alignment = ALIGN_MAP.get(align, PP_ALIGN.CENTER)

    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide


def rich_text_slide(prs, runs, bg_color=(0, 0, 0), x=1, y=3, w=14, h=3,
                    align="center", notes=None):
    """
    Create a single slide with rich text (multiple fonts/colors).

    Args:
        runs: list of dicts, each with:
            - text: str
            - size: int (pt)
            - color: tuple (r,g,b)
            - bold: bool (optional)
            - font_name: str (optional)
            - italic: bool (optional)

    Example:
        rich_text_slide(prs, [
            {"text": "WHOLESALE", "size": 72, "color": (255,192,0), "bold": True, "font_name": "Impact"},
            {"text": " is the future", "size": 72, "color": (255,255,255)},
        ])
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*bg_color)

    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]

    for run_spec in runs:
        run = p.add_run()
        run.text = run_spec["text"]
        run.font.size = Pt(run_spec.get("size", 40))
        run.font.color.rgb = RGBColor(*run_spec.get("color", (255, 255, 255)))
        run.font.bold = run_spec.get("bold", False)
        run.font.italic = run_spec.get("italic", False)
        if run_spec.get("font_name"):
            run.font.name = run_spec["font_name"]

    p.alignment = ALIGN_MAP.get(align, PP_ALIGN.CENTER)

    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide


def image_slide(prs, image_path, bg_color=(0, 0, 0), x=0, y=0, w=16, h=9,
                text=None, text_x=1, text_y=6, text_w=14, text_h=2,
                text_size=44, text_color=(255, 255, 255), text_bold=True,
                text_align="center", text_font=None,
                overlay=None, notes=None):
    """
    Create a slide with a background image, optional dark overlay, and optional text.

    Args:
        image_path: path to image file
        bg_color: fallback color if image fails
        x, y, w, h: image position and size (default: full slide)
        text: optional text to place over the image
        overlay: tuple (r, g, b, opacity) — e.g. (0, 0, 0, 0.6)
        notes: speaker notes

    Example (Fladlien-style atmospheric slide):
        image_slide(prs,
            "images/chess-pieces.png",
            overlay=(0, 0, 0, 0.55),
            text="The deck has been reshuffled...",
            text_size=56, text_bold=True,
            notes="Full speaker script here.")
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Solid background as fallback
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*bg_color)

    # Image
    pic = slide.shapes.add_picture(
        image_path,
        Inches(x), Inches(y),
        Inches(w), Inches(h)
    )

    # Move image to back if full-slide
    if x == 0 and y == 0:
        sp_tree = slide.shapes._spTree
        sp_tree.remove(pic._element)
        sp_tree.insert(2, pic._element)

    # Overlay
    if overlay:
        r, g, b, opacity = overlay
        _create_overlay(slide, r, g, b, opacity, Inches(w), Inches(h))

    # Text on top
    if text:
        box = slide.shapes.add_textbox(
            Inches(text_x), Inches(text_y),
            Inches(text_w), Inches(text_h)
        )
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(text_size)
        run.font.bold = text_bold
        run.font.color.rgb = RGBColor(*text_color)
        if text_font:
            run.font.name = text_font
        p.alignment = ALIGN_MAP.get(text_align, PP_ALIGN.CENTER)

    # Notes
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide


def split_slide(prs, image_path, text_runs, bg_color=(0, 0, 0),
                image_side="right", split_ratio=0.5,
                text_x=None, text_y=2, text_w=None, text_h=5,
                notes=None):
    """
    Create a two-column slide: text on one side, image on the other.

    Args:
        image_path: path to image file
        text_runs: list of rich text run dicts (same format as rich_text_slide)
        image_side: "left" or "right"
        split_ratio: 0.0-1.0 — portion of slide width for text (default 0.5)
        notes: speaker notes

    Example (Fladlien proof slide — headshot right, text left):
        split_slide(prs,
            "images/alex-hormozi.png",
            text_runs=[
                {"text": "He pays me $25,000.", "size": 48, "color": (255,192,0), "bold": True},
            ],
            image_side="right",
            notes="Full speaker script here.")
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*bg_color)

    slide_w = 16  # inches
    slide_h = 9
    text_portion = split_ratio
    img_portion = 1 - split_ratio

    if image_side == "right":
        # Text left, image right
        t_x = text_x if text_x is not None else 1
        t_w = text_w if text_w is not None else (slide_w * text_portion - 2)
        i_x = slide_w * text_portion
        i_w = slide_w * img_portion
    else:
        # Image left, text right
        i_x = 0
        i_w = slide_w * img_portion
        t_x = text_x if text_x is not None else (slide_w * img_portion + 1)
        t_w = text_w if text_w is not None else (slide_w * text_portion - 2)

    # Image
    slide.shapes.add_picture(
        image_path,
        Inches(i_x), Inches(0),
        Inches(i_w), Inches(slide_h)
    )

    # Text
    box = slide.shapes.add_textbox(
        Inches(t_x), Inches(text_y),
        Inches(t_w), Inches(text_h)
    )
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]

    for run_spec in text_runs:
        run = p.add_run()
        run.text = run_spec["text"]
        run.font.size = Pt(run_spec.get("size", 40))
        run.font.color.rgb = RGBColor(*run_spec.get("color", (255, 255, 255)))
        run.font.bold = run_spec.get("bold", False)
        run.font.italic = run_spec.get("italic", False)
        if run_spec.get("font_name"):
            run.font.name = run_spec["font_name"]

    p.alignment = PP_ALIGN.LEFT

    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide


def cta_footer(slide, url, urgency_text="Offer expires soon. Sign up now.",
               url_color=(0, 0, 255), text_color=(180, 180, 180),
               y=8.2, url_size=14, text_size=10):
    """
    Add persistent CTA footer to a slide (for offer section).

    Args:
        slide: python-pptx slide object
        url: the CTA URL string
        urgency_text: small text below URL
        url_color: tuple (r,g,b) — default blue
        text_color: tuple (r,g,b) — default light gray

    Example:
        slide = simple_slide(prs, "Bonus #1: The Oracle", ...)
        cta_footer(slide, "fasttrackformula.com")
    """
    # URL
    url_box = slide.shapes.add_textbox(
        Inches(1), Inches(y), Inches(14), Inches(0.4)
    )
    tf = url_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = url
    run.font.size = Pt(url_size)
    run.font.color.rgb = RGBColor(*url_color)
    run.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Urgency text
    urg_box = slide.shapes.add_textbox(
        Inches(1), Inches(y + 0.35), Inches(14), Inches(0.3)
    )
    tf2 = urg_box.text_frame
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = urgency_text
    run2.font.size = Pt(text_size)
    run2.font.color.rgb = RGBColor(*text_color)
    p2.alignment = PP_ALIGN.CENTER
